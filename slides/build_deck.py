from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "engineering-intelligence-board-deck.pptx"

SLIDES = [
    ("Engineering Intelligence Transformation Program", ["From reactive engineering to supervised autonomous systems", "18–24 month roadmap for knowledge, quality, resilience and self-healing infrastructure"]),
    ("The Scaling Problem", ["Knowledge fragmentation grows with team scale", "Senior engineers become bottlenecks", "Regression risk, MTTR and architecture drift rise with complexity"]),
    ("The Strategic Opportunity", ["Create a Private Engineering Intelligence Platform", "Ground AI in code, tickets, ADRs, incidents and telemetry", "Embed intelligence into IDE, PR, CI/CD and operations"]),
    ("What This Is — and Is Not", ["Secure cognitive infrastructure", "Not a chatbot rollout", "Not unrestricted production autonomy", "Not engineer replacement"]),
    ("Target Architecture", ["Entra ID + private networking", "Enterprise model gateway", "Hybrid retrieval with citations", "Azure DevOps, AKS, observability, policy and runbooks"]),
    ("Phase 1 — Institutional Memory", ["Continuous ingestion", "Citation-backed engineering Q&A", "IDE integration and faster onboarding"]),
    ("Phase 2 — Intelligent Guardrails", ["PR Guardian", "Security/architecture/IaC policy checks", "Historical regression matching"]),
    ("Phase 3 — Incident Intelligence", ["Correlate telemetry with deployment history", "Rank likely root causes", "Recommend rollback or runbook"]),
    ("Phase 4–5 — Predictive + Self-Healing", ["Deployment risk scoring", "Drift detection", "Policy-approved remediation", "Closed-loop verification"]),
    ("Financial Impact", ["Reclaimed engineering capacity", "Lower incident cost", "Controlled model and platform spend", "Stage-gated investment by ROI"]),
    ("Governance and Safety", ["RBAC before retrieval", "Human gates for high-risk actions", "Allow-listed deterministic runbooks", "Audit trail and kill switch"]),
    ("Strategic Outcome", ["Persistent engineering knowledge", "Measurable operational risk", "Faster delivery", "Supervised autonomous operations"]),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for title, bullets in SLIDES:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.8), Inches(1.0)).text_frame
    p = tx.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    body = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.2), Inches(4.8)).text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = bullet
        para.font.size = Pt(21)
        para.space_after = Pt(14)

prs.save(OUT)
print(OUT)
